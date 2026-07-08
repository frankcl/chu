"""Tests for skills/ppt/build.py — block normalization + .pptx rendering.

build.py is a standalone CLI script (not part of the agent package), so we load
it by path with importlib rather than importing it as a module.
"""

import importlib.util
from pathlib import Path

import pytest

_BUILD_PY = Path(__file__).resolve().parent.parent / "skills" / "ppt" / "build.py"
_THEMES_PY = Path(__file__).resolve().parent.parent / "skills" / "ppt" / "themes.py"


def _load_build():
    spec = importlib.util.spec_from_file_location("ppt_build", _BUILD_PY)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


def _load_themes():
    spec = importlib.util.spec_from_file_location("ppt_themes", _THEMES_PY)
    mod = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(mod)
    return mod


@pytest.fixture
def themes_mod():
    return _load_themes()


@pytest.fixture
def build_mod(tmp_path, monkeypatch):
    """build.py with its output dir redirected into tmp_path (no repo pollution)."""
    mod = _load_build()
    monkeypatch.setattr(mod, "_generated_dir", lambda: str(tmp_path))
    return mod


def _open(path):
    from pptx import Presentation
    return Presentation(path)


# ── _normalize_block ─────────────────────────────────────────────────────────

class TestNormalizeBlock:
    def test_bare_string_is_paragraph(self, build_mod):
        assert build_mod._normalize_block("hello") == {"type": "paragraph", "text": "hello"}

    def test_empty_string_dropped(self, build_mod):
        assert build_mod._normalize_block("   ") is None

    def test_explicit_paragraph(self, build_mod):
        out = build_mod._normalize_block({"type": "paragraph", "text": "a body"})
        assert out == {"type": "paragraph", "text": "a body"}

    def test_paragraph_alt_key(self, build_mod):
        out = build_mod._normalize_block({"type": "paragraph", "content": "via content"})
        assert out["type"] == "paragraph" and out["text"] == "via content"

    def test_bullets_block(self, build_mod):
        out = build_mod._normalize_block({"type": "bullets", "items": ["a", "b"]})
        assert out == {"type": "bullets", "items": ["a", "b"]}

    def test_table_block(self, build_mod):
        out = build_mod._normalize_block(
            {"type": "table", "headers": ["A", "B"], "rows": [["1", "2"], ["3", "4"]]})
        assert out["type"] == "table"
        assert out["headers"] == ["A", "B"]
        assert out["rows"] == [["1", "2"], ["3", "4"]]

    def test_table_alt_keys(self, build_mod):
        out = build_mod._normalize_block({"columns": ["X"], "data": [["1"], ["2"]]})
        assert out["type"] == "table"
        assert out["headers"] == ["X"]
        assert out["rows"] == [["1"], ["2"]]

    def test_table_rows_as_objects(self, build_mod):
        out = build_mod._normalize_block(
            {"type": "table", "headers": ["k", "v"],
             "rows": [{"k": "a", "v": "b"}, {"k": "c", "v": "d"}]})
        assert out["rows"] == [["a", "b"], ["c", "d"]]

    def test_numeric_cells_coerced_to_str(self, build_mod):
        out = build_mod._normalize_block({"type": "table", "rows": [[1, 2], [3, 4]]})
        assert out["rows"] == [["1", "2"], ["3", "4"]]

    def test_chart_dict_series(self, build_mod):
        out = build_mod._normalize_block(
            {"type": "chart", "chart_type": "line",
             "categories": ["Q1", "Q2"], "series": {"销量": [10, 20]}})
        assert out["type"] == "chart"
        assert out["chart_type"] == "line"
        assert out["categories"] == ["Q1", "Q2"]
        assert out["series"] == {"销量": [10.0, 20.0]}

    def test_chart_series_list_of_objects(self, build_mod):
        out = build_mod._normalize_block(
            {"type": "chart", "categories": ["A", "B"],
             "series": [{"name": "s1", "values": [1, 2]}, {"name": "s2", "values": [3, 4]}]})
        assert out["series"] == {"s1": [1.0, 2.0], "s2": [3.0, 4.0]}

    def test_chart_bare_number_series(self, build_mod):
        out = build_mod._normalize_block(
            {"type": "chart", "categories": ["A", "B"], "values": [5, 6]})
        assert out["series"] == {"系列1": [5.0, 6.0]}

    def test_chart_non_numeric_values_become_zero(self, build_mod):
        out = build_mod._normalize_block(
            {"type": "chart", "categories": ["A"], "series": {"s": ["x", 3]}})
        assert out["series"] == {"s": [0.0, 3.0]}

    def test_chart_without_data_dropped(self, build_mod):
        assert build_mod._normalize_block({"type": "chart", "categories": ["A"]}) is None

    def test_chart_kind_aliases(self, build_mod):
        assert build_mod._chart_kind({"chart_type": "饼图"}) == "pie"
        assert build_mod._chart_kind({"chart_type": "趋势"}) == "line"
        assert build_mod._chart_kind({"chart_type": "条形"}) == "bar"
        assert build_mod._chart_kind({}) == "column"

    def test_image_with_src(self, build_mod):
        out = build_mod._normalize_block(
            {"type": "image", "src": "https://x/y.png", "caption": "图1"})
        assert out == {"type": "image", "src": "https://x/y.png",
                       "prompt": None, "caption": "图1"}

    def test_image_with_prompt_only(self, build_mod):
        out = build_mod._normalize_block({"type": "image", "prompt": "a cat"})
        assert out["type"] == "image"
        assert out["src"] is None
        assert out["prompt"] == "a cat"

    def test_image_duck_typed_via_url(self, build_mod):
        out = build_mod._normalize_block({"url": "http://x/z.jpg"})
        assert out["type"] == "image"
        assert out["src"] == "http://x/z.jpg"

    def test_image_explicit_but_empty_dropped(self, build_mod):
        assert build_mod._normalize_block({"type": "image"}) is None


# ── _normalize_slide ─────────────────────────────────────────────────────────

class TestNormalizeSlide:
    def test_blocks_path(self, build_mod):
        out = build_mod._normalize_slide(
            {"title": "T", "blocks": [{"type": "paragraph", "text": "p"}]}, 1)
        assert out["title"] == "T"
        assert out["blocks"] == [{"type": "paragraph", "text": "p"}]

    def test_legacy_bullets_tagged(self, build_mod):
        out = build_mod._normalize_slide({"title": "T", "bullets": ["a", "b"]}, 1)
        assert len(out["blocks"]) == 1
        blk = out["blocks"][0]
        assert blk["type"] == "bullets"
        assert blk["items"] == ["a", "b"]
        assert blk["_legacy"] is True

    def test_string_slide(self, build_mod):
        out = build_mod._normalize_slide("just a title", 2)
        assert out["title"] == "just a title"
        assert out["blocks"] == []

    def test_missing_title_defaults(self, build_mod):
        out = build_mod._normalize_slide({"bullets": ["x"]}, 3)
        assert out["title"] == "第 3 页"


# ── build() end-to-end ───────────────────────────────────────────────────────

class TestBuild:
    def test_mixed_blocks_render(self, build_mod):
        spec = {
            "title": "Deck",
            "slides": [
                {"title": "Rich", "blocks": [
                    {"type": "paragraph", "text": "成段正文阐述"},
                    {"type": "bullets", "items": ["要点一", "要点二"]},
                    {"type": "table", "headers": ["车企", "销量"],
                     "rows": [["比亚迪", "300万"], ["特斯拉", "60万"]]},
                ]},
            ],
        }
        path = build_mod.build(spec)
        prs = _open(path)
        # title slide + 1 content slide
        assert len(prs.slides) == 2
        content = prs.slides[1]

        tables = [sh for sh in content.shapes if sh.has_table]
        assert len(tables) == 1
        tbl = tables[0].table
        assert tbl.cell(0, 0).text == "车企"
        assert tbl.cell(1, 0).text == "比亚迪"

        all_text = "\n".join(
            sh.text_frame.text for sh in content.shapes if sh.has_text_frame)
        assert "成段正文阐述" in all_text
        assert "要点一" in all_text

    def test_legacy_bullets_still_build(self, build_mod):
        spec = {"slides": [{"title": "Plain", "bullets": ["a", "b", "c"]}]}
        path = build_mod.build(spec)
        prs = _open(path)
        assert len(prs.slides) == 2
        text = "\n".join(
            sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame)
        assert "a" in text and "Plain" in text

    def test_malformed_table_block_does_not_crash(self, build_mod):
        # rows is not a list → block should be dropped during normalization,
        # and even if a degenerate block slips through, build must not raise.
        spec = {"slides": [{"title": "Bad", "blocks": [
            {"type": "table", "headers": ["A"], "rows": "not-a-list"},
            {"type": "paragraph", "text": "survives"},
        ]}]}
        path = build_mod.build(spec)
        prs = _open(path)
        assert len(prs.slides) == 2
        text = "\n".join(
            sh.text_frame.text for sh in prs.slides[1].shapes if sh.has_text_frame)
        assert "survives" in text

    def test_notes_attached(self, build_mod):
        spec = {"slides": [{"title": "N", "blocks": [
            {"type": "paragraph", "text": "body"}], "notes": "speaker note"}]}
        path = build_mod.build(spec)
        prs = _open(path)
        assert prs.slides[1].notes_slide.notes_text_frame.text == "speaker note"

    def test_dark_theme_table_cells_get_explicit_fill(self, build_mod):
        """On a dark theme, table body cells must carry an explicit fill so the
        light body text isn't drawn onto python-pptx's default light table style."""
        from pptx.util import Emu  # noqa: F401 — ensure pptx importable
        spec = {"theme": "tech-dark", "slides": [{"title": "Data", "blocks": [
            {"type": "table", "headers": ["A", "B"], "rows": [["1", "2"]]},
        ]}]}
        path = build_mod.build(spec)
        prs = _open(path)
        tables = [sh.table for sh in prs.slides[1].shapes if sh.has_table]
        assert tables, "expected a table on the slide"
        body_cell = tables[0].cell(1, 0)
        assert str(body_cell.fill.fore_color.rgb) == "16213E"


# ── themes module (shared palette / preview source of truth) ─────────────────

class TestThemesModule:
    def test_build_uses_shared_themes(self, build_mod, themes_mod):
        """build.py renders from the same THEMES dict the server serves to the UI."""
        assert build_mod._THEMES is themes_mod.THEMES or build_mod._THEMES == themes_mod.THEMES
        assert themes_mod.THEMES["tech-dark"]["band"] == "2E6CB5"

    def test_theme_previews_shape(self, themes_mod):
        previews = themes_mod.theme_previews()
        assert [p["name"] for p in previews] == [
            "default", "business-blue", "tech-dark", "minimal"]
        for p in previews:
            assert p["label"]
            assert set(p["colors"]) >= {
                "background", "bg2", "band", "band_text", "text", "title", "footer"}
            assert p["sample"]["title"] and len(p["sample"]["body"]) >= 1

    def test_preview_colors_match_palette(self, themes_mod):
        biz = next(p for p in themes_mod.theme_previews() if p["name"] == "business-blue")
        assert biz["colors"]["band"] == "1F4E79"
        assert biz["colors"]["background"] == "FFFFFF"
        # the plain `default` theme has no band/background (renders as a blank slide)
        dft = next(p for p in themes_mod.theme_previews() if p["name"] == "default")
        assert dft["colors"]["band"] is None and dft["colors"]["background"] is None


# ── theme resolution ─────────────────────────────────────────────────────────

class TestResolveTheme:
    def test_named_theme(self, build_mod):
        t = build_mod._resolve_theme({"theme": "tech-dark"})
        assert t["background"] == "1A1A2E"
        assert t["bg2"] == "16213E"        # gradient end
        assert t["band"] == "2E6CB5"       # header band color (brighter than bg)
        assert t["table_body"] == "16213E"  # explicit dark cell fill for readable text

    def test_band_distinct_from_background(self, build_mod):
        """Header band / footer must not collapse into the slide background."""
        for name in ("business-blue", "tech-dark", "minimal"):
            t = build_mod._THEMES[name]
            assert t["band"] not in (t["background"], t["bg2"])
            assert t["footer"] not in (t["background"], t["bg2"])

    def test_unknown_name_falls_back_to_default(self, build_mod):
        t = build_mod._resolve_theme({"theme": "nope"})
        assert t == build_mod._THEMES["default"]

    def test_no_theme_is_default(self, build_mod):
        assert build_mod._resolve_theme({}) == build_mod._THEMES["default"]

    def test_inline_overrides_named_base(self, build_mod):
        t = build_mod._resolve_theme(
            {"theme": {"base": "minimal", "title": "1F4E79", "table_header": "2E75B6"}})
        assert t["title"] == "1F4E79"           # overridden
        assert t["table_header"] == "2E75B6"    # overridden
        assert t["text"] == build_mod._THEMES["minimal"]["text"]  # inherited

    def test_style_alias(self, build_mod):
        # `style` is accepted as an alias for `theme`.
        assert build_mod._resolve_theme({"style": "minimal"})["title"] == "111111"


# ── _parse_spec ──────────────────────────────────────────────────────────────

class TestParseSpec:
    def test_plain_object(self, build_mod):
        spec = build_mod._parse_spec('{"title": "T", "slides": []}')
        assert spec["title"] == "T"

    def test_literal_newline_in_string_tolerated(self, build_mod):
        """A real newline inside a string value (the model's `\\n` decoded once
        before reaching build.py) must not fail parsing — the prior strict parse
        raised 'Invalid control character' and aborted the whole build."""
        raw = '{"slides": [{"title": "S", "blocks": [{"type": "paragraph", "text": "a\nb"}]}]}'
        assert "\n" in raw  # genuinely a control character, not an escape
        spec = build_mod._parse_spec(raw)
        assert spec["slides"][0]["blocks"][0]["text"] == "a\nb"

    def test_double_encoded(self, build_mod):
        import json
        inner = json.dumps({"title": "T", "slides": []})
        spec = build_mod._parse_spec(json.dumps(inner))  # encoded twice
        assert spec["title"] == "T"

    def test_single_element_array_unwrapped(self, build_mod):
        spec = build_mod._parse_spec('[{"title": "T", "slides": []}]')
        assert spec["title"] == "T"

    def test_bare_slide_list_wrapped(self, build_mod):
        spec = build_mod._parse_spec('[{"title": "S", "blocks": []}]')
        assert "slides" in spec and len(spec["slides"]) == 1


# ── main(): argv handling ─────────────────────────────────────────────────────

class TestMainArgv:
    _SPEC = ('{"title": "T", "slides": [{"title": "S", "blocks": ['
             '{"type": "paragraph", "text": "hello world"}]}]}')

    def _run_main(self, build_mod, monkeypatch, capsys, argv_tail):
        import json
        import sys
        monkeypatch.setattr(sys, "argv", ["build.py", *argv_tail])
        build_mod.main()  # prints structured JSON; no sys.exit on success
        return json.loads(capsys.readouterr().out)

    def test_single_arg_builds(self, build_mod, monkeypatch, capsys):
        out = self._run_main(build_mod, monkeypatch, capsys, [self._SPEC])
        assert out["ok"] is True and out["slides"] == 1

    def test_split_across_multiple_args_is_rejoined(self, build_mod, monkeypatch, capsys):
        """The model sometimes passes the spec split across script_args elements;
        each arrives as a separate argv. main() must rejoin them (previously it
        read only argv[1] and failed with 'Expecting , delimiter')."""
        # split at an arbitrary mid-string boundary (mimics char-434 truncation)
        a, b, c = self._SPEC[:30], self._SPEC[30:55], self._SPEC[55:]
        out = self._run_main(build_mod, monkeypatch, capsys, [a, b, c])
        assert out["ok"] is True and out["slides"] == 1

    def test_empty_args_usage_error(self, build_mod, monkeypatch, capsys):
        import pytest
        with pytest.raises(SystemExit):
            self._run_main(build_mod, monkeypatch, capsys, ["   "])
        assert "usage" in capsys.readouterr().out


class TestHexParsing:
    def test_six_digit(self, build_mod):
        assert str(build_mod._hex("1F4E79")) == "1F4E79"

    def test_hash_and_three_digit(self, build_mod):
        assert str(build_mod._hex("#abc")) == "AABBCC"

    def test_invalid_returns_none(self, build_mod):
        assert build_mod._hex("not-a-color") is None
        assert build_mod._hex(None) is None


# ── layout picking ───────────────────────────────────────────────────────────

class TestPickLayout:
    def test_matches_by_name(self, build_mod):
        from pptx import Presentation
        prs = Presentation()
        layout = build_mod._pick_layout(prs, ["title only"], 99)
        assert "title only" in layout.name.lower()

    def test_falls_back_to_index_when_no_name_match(self, build_mod):
        from pptx import Presentation
        prs = Presentation()
        layout = build_mod._pick_layout(prs, ["nonexistent-layout"], 1)
        assert layout is list(prs.slide_layouts)[1]


# ── theme applied end-to-end ─────────────────────────────────────────────────

class TestThemedBuild:
    def test_themed_build_succeeds(self, build_mod):
        spec = {
            "theme": "tech-dark",
            "slides": [{"title": "T", "blocks": [
                {"type": "paragraph", "text": "正文"},
                {"type": "table", "headers": ["A"], "rows": [["1"]]},
            ]}],
        }
        path = build_mod.build(spec)
        prs = _open(path)
        assert len(prs.slides) == 2
        # table header cell uses the theme's header fill color
        from pptx.dml.color import RGBColor
        tbl = next(sh.table for sh in prs.slides[1].shapes if sh.has_table)
        assert tbl.cell(0, 0).fill.fore_color.rgb == RGBColor.from_string("2E6CB5")

    def test_default_theme_keeps_legacy_header_color(self, build_mod):
        spec = {"slides": [{"title": "T", "blocks": [
            {"type": "table", "headers": ["A"], "rows": [["1"]]}]}]}
        path = build_mod.build(spec)
        prs = _open(path)
        from pptx.dml.color import RGBColor
        tbl = next(sh.table for sh in prs.slides[1].shapes if sh.has_table)
        assert tbl.cell(0, 0).fill.fore_color.rgb == RGBColor.from_string("1F4E79")


# ── background / decoration (header band, gradient, footer) ───────────────────

def _auto_shapes(slide):
    """Decorative AUTO_SHAPE rectangles added by the background decoration."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    return [sh for sh in slide.shapes if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE]


class TestBackgroundDecoration:
    def test_default_theme_has_no_decoration(self, build_mod):
        """Zero regression: default deck gets no band/footer shapes."""
        spec = {"slides": [{"title": "T", "blocks": [
            {"type": "paragraph", "text": "x"}]}]}
        prs = _open(build_mod.build(spec))
        assert _auto_shapes(prs.slides[1]) == []
        assert _auto_shapes(prs.slides[0]) == []

    def test_named_theme_adds_header_band_with_white_title(self, build_mod):
        from pptx.dml.color import RGBColor
        spec = {"theme": "business-blue",
                "slides": [{"title": "标题", "blocks": [{"type": "paragraph", "text": "x"}]}]}
        prs = _open(build_mod.build(spec))
        content = prs.slides[1]
        # at least one decorative rectangle (header band + footer bar)
        assert len(_auto_shapes(content)) >= 1
        # title text painted in band_text (white) for contrast on the band
        title = content.shapes.title
        assert title.text_frame.paragraphs[0].font.color.rgb == RGBColor.from_string("FFFFFF")

    def test_minimal_band_uses_white_title_text(self, build_mod):
        # minimal now uses a solid dark-gray band (so it's visible on white),
        # with white title text painted on it for contrast.
        from pptx.dml.color import RGBColor
        spec = {"theme": "minimal",
                "slides": [{"title": "标题", "blocks": [{"type": "paragraph", "text": "x"}]}]}
        prs = _open(build_mod.build(spec))
        title = prs.slides[1].shapes.title
        assert title.text_frame.paragraphs[0].font.color.rgb == RGBColor.from_string("FFFFFF")

    def test_gradient_background_when_bg2_set(self, build_mod):
        from pptx.enum.dml import MSO_FILL
        spec = {"theme": "business-blue", "slides": [{"title": "T", "bullets": ["a"]}]}
        prs = _open(build_mod.build(spec))
        assert prs.slides[0].background.fill.type == MSO_FILL.GRADIENT

    def test_solid_background_when_no_bg2(self, build_mod):
        from pptx.enum.dml import MSO_FILL
        spec = {"theme": "minimal", "slides": [{"title": "T", "bullets": ["a"]}]}
        prs = _open(build_mod.build(spec))
        assert prs.slides[0].background.fill.type == MSO_FILL.SOLID

    def test_header_band_is_behind_title(self, build_mod):
        """The band must sit below the title in z-order (earlier in spTree)."""
        spec = {"theme": "tech-dark",
                "slides": [{"title": "标题", "blocks": [{"type": "paragraph", "text": "x"}]}]}
        prs = _open(build_mod.build(spec))
        content = prs.slides[1]
        sp_tree = content.shapes._spTree
        order = list(sp_tree)
        band = _auto_shapes(content)[0]._element
        title_el = content.shapes.title._element
        assert order.index(band) < order.index(title_el)

    def test_footer_page_number_present(self, build_mod):
        spec = {"theme": "business-blue", "slides": [
            {"title": "P1", "bullets": ["a"]},
            {"title": "P2", "bullets": ["b"]},
        ]}
        prs = _open(build_mod.build(spec))
        # slides[2] is the 2nd content slide → page number "2"
        texts = [sh.text_frame.text for sh in prs.slides[2].shapes if sh.has_text_frame]
        assert "2" in texts

    def test_themed_blocks_legacy_cover_all_build(self, build_mod):
        """A banded theme builds cleanly across cover + blocks + legacy pages."""
        spec = {"theme": "tech-dark", "subtitle": "副标题", "slides": [
            {"title": "Blocks", "blocks": [
                {"type": "paragraph", "text": "正文"},
                {"type": "table", "headers": ["A"], "rows": [["1"]]}]},
            {"title": "Legacy", "bullets": ["a", "b"]},
        ]}
        prs = _open(build_mod.build(spec))
        assert len(prs.slides) == 3  # cover + 2 content


# ── title fit inside the band (vertical centering + overflow) ────────────────

class TestTitleFit:
    def test_font_size_shrinks_for_longer_titles(self, build_mod):
        short = build_mod._title_font_size("短标题")
        medium = build_mod._title_font_size("这是一个比较长一些的幻灯片标题示例")
        long = build_mod._title_font_size("这是一个非常非常长的幻灯片标题" * 3)
        assert short > medium >= long
        assert long >= 18  # never absurdly small

    def test_banded_title_is_vertically_centered_and_wrapped(self, build_mod):
        from pptx.enum.text import MSO_ANCHOR
        from pptx.util import Inches
        spec = {"theme": "business-blue", "slides": [
            {"title": "一个相当长的标题用于触发换行以验证不会溢出色块背景区域", "bullets": ["a"]}]}
        prs = _open(build_mod.build(spec))
        title = prs.slides[1].shapes.title
        tf = title.text_frame
        assert tf.vertical_anchor == MSO_ANCHOR.MIDDLE
        assert tf.word_wrap is True
        # title repositioned to the top band region, height == band height
        assert title.top == 0
        assert abs(title.height - Inches(build_mod._BAND_HEIGHT)) < Inches(0.02)

    def test_default_theme_title_not_repositioned(self, build_mod):
        # No band → title keeps the layout's placeholder geometry (top != 0).
        spec = {"slides": [{"title": "标题", "bullets": ["a"]}]}
        prs = _open(build_mod.build(spec))
        assert prs.slides[1].shapes.title.top != 0


# ── webp (and other non-native formats) image support ────────────────────────

class TestWebpImages:
    def _make_webp(self, tmp_path):
        from PIL import Image
        p = tmp_path / "pic.webp"
        Image.new("RGB", (4, 4), (200, 30, 30)).save(str(p), "WEBP")
        return p

    def test_ensure_supported_converts_webp_to_png(self, build_mod, tmp_path):
        webp = self._make_webp(tmp_path)
        out = build_mod._ensure_supported_image(str(webp))
        assert out.lower().endswith(".png")
        from pptx import Presentation  # the converted file must be embeddable
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        from pptx.util import Inches
        slide.shapes.add_picture(out, Inches(1), Inches(1), width=Inches(1))

    def test_native_ext_passthrough(self, build_mod, tmp_path):
        png = tmp_path / "x.png"
        png.write_bytes(_PNG_1x1)
        assert build_mod._ensure_supported_image(str(png)) == str(png)

    def test_webp_image_block_renders(self, build_mod, tmp_path):
        webp = self._make_webp(tmp_path)
        spec = {"slides": [{"title": "I", "blocks": [
            {"type": "image", "src": str(webp), "caption": "webp 图"}]}]}
        prs = _open(build_mod.build(spec))
        assert any(sh.shape_type == 13 for sh in prs.slides[1].shapes)  # PICTURE present


# ── cover header band (must not be thinner than the bottom accent) ───────────

class TestCoverHeader:
    def test_cover_top_band_not_thinner_than_bottom(self, build_mod):
        spec = {"theme": "business-blue", "slides": [{"title": "T", "bullets": ["a"]}]}
        prs = _open(build_mod.build(spec))
        shapes = _auto_shapes(prs.slides[0])
        assert len(shapes) >= 2
        top_shape = min(shapes, key=lambda s: s.top)     # the header (top == 0)
        bottom_shape = max(shapes, key=lambda s: s.top)  # the bottom accent
        assert top_shape.top == 0
        assert top_shape.height >= bottom_shape.height   # header at least as tall


# ── content must not overflow into the footer ────────────────────────────────

class TestContentOverflow:
    def test_tall_blocks_dropped_above_footer(self, build_mod):
        from pptx.util import Inches
        charts = [{"type": "chart", "chart_type": "column",
                   "categories": ["a", "b"], "series": {"s": [1, 2]}} for _ in range(3)]
        spec = {"slides": [{"title": "T", "blocks": charts}]}
        prs = _open(build_mod.build(spec))
        placed = [sh for sh in prs.slides[1].shapes if sh.has_chart]
        # only what fits above the footer is drawn; the rest are dropped
        assert len(placed) == 1
        assert placed[0].top + placed[0].height <= Inches(7.0)

    def test_first_block_always_placed_even_if_tall(self, build_mod):
        # A single very tall block is still drawn (a page is never left empty).
        spec = {"slides": [{"title": "T", "blocks": [
            {"type": "chart", "chart_type": "column",
             "categories": ["a"], "series": {"s": [1]}}]}]}
        prs = _open(build_mod.build(spec))
        assert any(sh.has_chart for sh in prs.slides[1].shapes)


# ── image resilience: one failure must not drop the others ───────────────────

class TestImageResilience:
    def test_budget_skips_network_but_keeps_local(self, build_mod, tmp_path, monkeypatch):
        local = tmp_path / "ok.png"
        local.write_bytes(_PNG_1x1)
        monkeypatch.setattr(build_mod, "_img_time_spent", build_mod._IMG_TIME_BUDGET + 1)
        # network fetch skipped once the budget is spent …
        assert build_mod._resolve_image("https://x/y.png", None) is None
        # … but a local file always resolves regardless of the budget.
        assert build_mod._resolve_image(str(local), None) == str(local)

    def test_one_failed_download_keeps_remaining_images(self, build_mod, tmp_path, monkeypatch):
        good = tmp_path / "good.png"
        good.write_bytes(_PNG_1x1)
        # Simulate the first (URL) image failing to download; the second is local.
        monkeypatch.setattr(build_mod, "_download_image", lambda url: None)
        spec = {"slides": [{"title": "I", "blocks": [
            {"type": "image", "src": "https://x/bad.png"},
            {"type": "image", "src": str(good)},
        ]}]}
        prs = _open(build_mod.build(spec))
        pics = [sh for sh in prs.slides[1].shapes if sh.shape_type == 13]
        assert len(pics) == 1  # the good image still renders despite the failed one


# ── template loading ─────────────────────────────────────────────────────────

class TestTemplateLoading:
    def test_missing_template_falls_back(self, build_mod):
        # A non-existent template name must not crash — it falls back to default.
        spec = {"template": "does-not-exist", "slides": [{"title": "T", "bullets": ["a"]}]}
        path = build_mod.build(spec)
        assert _open(path) is not None

    def test_path_traversal_template_ignored(self, build_mod):
        spec = {"template": "../../../etc/passwd", "slides": [{"title": "T", "bullets": ["a"]}]}
        path = build_mod.build(spec)  # must not raise / must not load outside templates/
        assert len(_open(path).slides) == 2

    def test_valid_template_is_used(self, build_mod, tmp_path, monkeypatch):
        # Create a real template file with a marker slide, point the loader at it,
        # and confirm build() starts from that template (marker slide preserved).
        from pptx import Presentation
        tdir = tmp_path / "templates"
        tdir.mkdir()
        base = Presentation()
        base.slides.add_slide(base.slide_layouts[6])  # a blank marker slide
        base.save(str(tdir / "mark.pptx"))
        monkeypatch.setattr(build_mod, "_templates_dir", lambda: str(tdir))

        spec = {"template": "mark", "slides": [{"title": "T", "bullets": ["a"]}]}
        prs = _open(build_mod.build(spec))
        # 1 marker (from template) + 1 title + 1 content
        assert len(prs.slides) == 3


# ── chart & image rendering ──────────────────────────────────────────────────

# A 1x1 PNG (smallest valid), base64-decoded at use sites.
_PNG_1x1 = bytes.fromhex(
    "89504e470d0a1a0a0000000d49484452000000010000000108060000001f15c4"
    "890000000d49444154789c6360000002000100ffff03000006000557bfabd400"
    "00000049454e44ae426082"
)


class TestChartRendering:
    def test_chart_renders_as_native_chart(self, build_mod):
        spec = {"slides": [{"title": "C", "blocks": [
            {"type": "chart", "chart_type": "column",
             "categories": ["Q1", "Q2", "Q3"], "series": {"销量": [80, 95, 110]}},
        ]}]}
        prs = _open(build_mod.build(spec))
        charts = [sh for sh in prs.slides[1].shapes if sh.has_chart]
        assert len(charts) == 1
        plot = charts[0].chart.plots[0]
        assert list(plot.categories) == ["Q1", "Q2", "Q3"]
        assert plot.series[0].values == (80.0, 95.0, 110.0)

    def test_pie_chart_single_series(self, build_mod):
        spec = {"slides": [{"title": "P", "blocks": [
            {"type": "chart", "chart_type": "pie",
             "categories": ["A", "B"], "series": {"s1": [60, 40], "s2": [1, 2]}},
        ]}]}
        prs = _open(build_mod.build(spec))
        chart = next(sh.chart for sh in prs.slides[1].shapes if sh.has_chart)
        # pie uses just the first series
        assert len(chart.plots[0].series) == 1


class TestImageRendering:
    def test_local_image_renders(self, build_mod, tmp_path):
        img = tmp_path / "pic.png"
        img.write_bytes(_PNG_1x1)
        spec = {"slides": [{"title": "I", "blocks": [
            {"type": "image", "src": str(img), "caption": "图注"},
        ]}]}
        prs = _open(build_mod.build(spec))
        pics = [sh for sh in prs.slides[1].shapes if sh.shape_type == 13]  # PICTURE
        assert len(pics) == 1
        text = "\n".join(sh.text_frame.text for sh in prs.slides[1].shapes
                         if sh.has_text_frame)
        assert "图注" in text

    def test_image_from_generated_dir_by_basename(self, build_mod, tmp_path):
        # _generated_dir is monkeypatched to tmp_path by the fixture.
        (tmp_path / "logo.png").write_bytes(_PNG_1x1)
        spec = {"slides": [{"title": "I", "blocks": [
            {"type": "image", "src": "logo.png"}]}]}
        prs = _open(build_mod.build(spec))
        assert any(sh.shape_type == 13 for sh in prs.slides[1].shapes)

    def test_unresolvable_image_skipped_not_fatal(self, build_mod):
        spec = {"slides": [{"title": "I", "blocks": [
            {"type": "image", "src": "/nope/missing.png"},
            {"type": "paragraph", "text": "survives"},
        ]}]}
        prs = _open(build_mod.build(spec))
        s = prs.slides[1]
        assert not any(sh.shape_type == 13 for sh in s.shapes)  # no picture
        text = "\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        assert "survives" in text

    def test_download_url_mocked(self, build_mod, tmp_path, monkeypatch):
        # Avoid real network: stub _download_image to drop a local file.
        target = tmp_path / "dl.png"
        target.write_bytes(_PNG_1x1)
        monkeypatch.setattr(build_mod, "_download_image", lambda url: str(target))
        spec = {"slides": [{"title": "I", "blocks": [
            {"type": "image", "src": "https://example.com/x.png"}]}]}
        prs = _open(build_mod.build(spec))
        assert any(sh.shape_type == 13 for sh in prs.slides[1].shapes)

    def test_prompt_generation_no_key_skips(self, build_mod, monkeypatch):
        monkeypatch.delenv("DASHSCOPE_API_KEY", raising=False)
        # prompt-only image with no backend → resolves to None, block skipped.
        assert build_mod._resolve_image(None, "a prompt") is None

    def test_image_survives_on_text_heavy_slide(self, build_mod, tmp_path):
        """Regression: an image appended after a full page of text used to be
        dropped by the vertical flow. Now text+image lay out side by side."""
        img = tmp_path / "pic.png"
        img.write_bytes(_PNG_1x1)
        para = "美伊关系长期处于紧张状态。" * 12  # long enough to fill the page vertically
        spec = {"theme": "tech-dark", "slides": [{"title": "背景", "blocks": [
            {"type": "paragraph", "text": para},
            {"type": "bullets", "items": ["要点一", "要点二", "要点三", "要点四", "要点五"]},
            {"type": "image", "src": str(img), "caption": "配图"},
        ]}]}
        prs = _open(build_mod.build(spec))
        s = prs.slides[1]
        assert any(sh.shape_type == 13 for sh in s.shapes), "image must be placed"
        text = "\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        assert "要点一" in text and "配图" in text  # text column + caption both present

    def test_image_with_table_still_renders_via_flow(self, build_mod, tmp_path):
        """image + a full-width table: flow layout shrinks the image to the room
        left instead of dropping it."""
        img = tmp_path / "pic.png"
        img.write_bytes(_PNG_1x1)
        spec = {"slides": [{"title": "数据", "blocks": [
            {"type": "table", "headers": ["A", "B"], "rows": [["1", "2"], ["3", "4"]]},
            {"type": "image", "src": str(img)},
        ]}]}
        prs = _open(build_mod.build(spec))
        assert any(sh.shape_type == 13 for sh in prs.slides[1].shapes)

    def test_two_column_falls_back_when_image_unresolvable(self, build_mod):
        """If the only image can't be resolved, the text must still render (the
        two-column path reports failure and the flow layout takes over)."""
        spec = {"slides": [{"title": "X", "blocks": [
            {"type": "paragraph", "text": "重要内容"},
            {"type": "image", "src": "/nope/missing.png"},
        ]}]}
        prs = _open(build_mod.build(spec))
        s = prs.slides[1]
        assert not any(sh.shape_type == 13 for sh in s.shapes)
        text = "\n".join(sh.text_frame.text for sh in s.shapes if sh.has_text_frame)
        assert "重要内容" in text

    def _real_png(self, tmp_path, w, h, name="p.png"):
        from PIL import Image
        p = tmp_path / name
        Image.new("RGB", (w, h), (80, 110, 160)).save(str(p))
        return str(p)

    def test_two_column_text_clears_image_no_overlap(self, build_mod, tmp_path):
        """The body text column must end to the left of the image (issue #1:
        text used to overflow its column and overlap the picture)."""
        img = self._real_png(tmp_path, 1600, 1000)
        para = "美伊关系长期紧张，核心争议围绕伊朗核问题展开。" * 4
        spec = {"theme": "business-blue", "slides": [{"title": "S", "blocks": [
            {"type": "paragraph", "text": para},
            {"type": "bullets", "items": [
                "伊朗外长表示谈判有好的开始但建立信任需要时间且充满变数与挑战",
                "美国总统称会谈进行得非常好并表示将再次安排新一轮的谈判",
            ]},
            {"type": "image", "src": img, "caption": "配图"},
        ]}]}
        prs = _open(build_mod.build(spec))
        s = prs.slides[1]
        pic = next(sh for sh in s.shapes if sh.shape_type == 13)
        body = next(sh for sh in s.shapes
                    if sh.has_text_frame and "美伊关系" in sh.text_frame.text)
        # text's right edge must not cross into the image (small EMU tolerance)
        assert body.left + body.width <= pic.left + 5000
        # and the body text wrapped inside a narrow column (< half the 10in slide)
        assert body.width < build_mod._EMU_PER_INCH * 5

    def test_two_column_image_is_large_and_not_top_hugging(self, build_mod, tmp_path):
        """Issue #2: the image should fill the right column width and be centered
        vertically (not pinned to the top with whitespace below)."""
        img = self._real_png(tmp_path, 1600, 1000)
        spec = {"theme": "tech-dark", "slides": [{"title": "S", "blocks": [
            {"type": "paragraph", "text": "简短说明。"},
            {"type": "image", "src": img, "caption": "图注"},
        ]}]}
        prs = _open(build_mod.build(spec))
        pic = next(sh for sh in prs.slides[1].shapes if sh.shape_type == 13)
        inch = build_mod._EMU_PER_INCH
        assert pic.width / inch >= 3.5            # fills (most of) the 4.0in column
        assert pic.top / inch > 1.9               # pushed below content top → centered
