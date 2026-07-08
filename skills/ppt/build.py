"""Build a PowerPoint (.pptx) from a JSON spec passed as the first CLI argument.

Usage: python build.py '<json-spec>'

Spec shape (see SKILL.md for details):
    {
      "title": "Deck title",
      "subtitle": "optional",
      "output": "deck.pptx",
      "theme": "business-blue",   # optional: default|business-blue|tech-dark|minimal,
                                  #   or an inline object {"base":..,"title":"1F4E79",..}
      "template": "my-template",  # optional: a .pptx in skills/ppt/templates/
      "slides": [
        {"title": "Heading", "bullets": ["a", "b"], "notes": "optional"},
        {"title": "Rich slide", "blocks": [
          {"type": "paragraph", "text": "成段正文……"},
          {"type": "bullets", "items": ["a", "b"]},
          {"type": "table", "headers": ["A", "B"], "rows": [["1", "2"]]}
        ]},
        ...
      ]
    }

A slide may carry either legacy `bullets` (a list of strings / nested bullet
objects) OR an ordered `blocks` list mixing block types:
  - {"type": "paragraph", "text": "..."}
  - {"type": "bullets", "items": [...]}
  - {"type": "table", "headers": [...], "rows": [[...], ...]}
  - {"type": "chart", "chart_type": "bar|column|line|pie",
     "categories": [...], "series": {"name": [numbers]}}
  - {"type": "image", "src": "<url|local path>", "caption": "...",
     "prompt": "<optional: generate via DashScope if no src>"}
Each bullet is either a string, or {"text": "...", "level": 0..2} for indenting.

Prints JSON: {"ok": true, "path": "<abs path>", "slides": N}
          or {"ok": false, "error": "<message>"}

Requires python-pptx (declared in the project's dependencies).
"""

import hashlib
import json
import os
import sys
import time
import urllib.request
import uuid


def _generated_dir() -> str:
    """Repo-root `generated/` dir — where decks are saved for the server to serve.

    build.py lives at <repo>/skills/ppt/build.py, so the repo root is two levels
    up from this file's directory.
    """
    here = os.path.dirname(os.path.abspath(__file__))
    repo_root = os.path.dirname(os.path.dirname(here))
    return os.path.join(repo_root, "generated")


# Keys a bullet object may use for its text / its nested sub-bullets. Models
# emit many shapes (e.g. {"point": "..."} or {"heading": "..", "points": [..]}),
# so we look these up rather than assuming a single "text" key.
_BULLET_TEXT_KEYS = ("text", "content", "point", "bullet", "label", "value", "name", "title", "heading")
_BULLET_CHILD_KEYS = ("children", "sub", "subs", "sub_bullets", "subbullets", "points", "bullets", "items", "details")


def _bullet_text(d: dict) -> str:
    for k in _BULLET_TEXT_KEYS:
        v = d.get(k)
        if isinstance(v, str) and v.strip():
            return v
        if isinstance(v, (int, float)):
            return str(v)
    return ""


def _bullet_children(d: dict):
    for k in _BULLET_CHILD_KEYS:
        v = d.get(k)
        if isinstance(v, list):
            return v
    return []


def _add_bullets(body_frame, bullets) -> None:
    """Fill a body placeholder with (possibly nested) bullets, shape-tolerant.

    Accepts each bullet as a string, a number, a nested list, or a dict using
    any of the common text / child keys. Empty parents are skipped (so their
    children still render) rather than emitting blank lines.
    """
    state = {"first": True}

    def emit(items, level: int) -> None:
        for b in items:
            children = []
            if isinstance(b, dict):
                text = _bullet_text(b)
                lvl = int(b.get("level", level) or level)
                children = _bullet_children(b)
            elif isinstance(b, (list, tuple)):
                emit(b, level)
                continue
            else:
                text = str(b)
                lvl = level
            if text:
                # Reuse the pre-existing first paragraph; add new ones after.
                para = body_frame.paragraphs[0] if state["first"] else body_frame.add_paragraph()
                para.text = text
                para.level = max(0, min(lvl, 4))  # python-pptx supports levels 0–4
                state["first"] = False
            if children:
                emit(children, min((lvl if text else level) + 1, 4))

    emit(bullets, 0)


# Alternative key names the model commonly emits, mapped to our canonical ones.
_TITLE_KEYS = ("title", "heading", "header", "name", "slide_title", "subject")
_BULLETS_KEYS = ("bullets", "content", "points", "items", "body", "text", "lines", "details")
_NOTES_KEYS = ("notes", "speaker_notes", "note")
_BLOCKS_KEYS = ("blocks", "sections", "elements")

# Per-block alternative keys (shape-tolerant, mirroring the bullet/title sets).
_PARA_TEXT_KEYS = ("text", "content", "paragraph", "body", "value")
_TABLE_HEADER_KEYS = ("headers", "header", "columns", "cols", "head")
_TABLE_ROWS_KEYS = ("rows", "data", "records", "body")
_IMAGE_TYPES = ("image", "picture", "img", "photo")
_IMAGE_SRC_KEYS = ("src", "url", "image", "path", "file", "uri")
_IMAGE_CAPTION_KEYS = ("caption", "alt", "label", "title")
_CHART_TYPES = ("chart", "graph", "plot")
_CHART_CAT_KEYS = ("categories", "labels", "x", "axis")
_CHART_SERIES_KEYS = ("series", "values", "data", "y")


def _first_key(d: dict, keys) -> object:
    for k in keys:
        if d.get(k):
            return d[k]
    return None


def _floats(values) -> list:
    """Coerce a sequence to a list of floats; non-numeric entries become 0.0."""
    out = []
    for v in values if isinstance(values, (list, tuple)) else []:
        try:
            out.append(float(v))
        except (TypeError, ValueError):
            out.append(0.0)
    return out


def _chart_kind(b: dict) -> str:
    """Normalize a requested chart kind to one of: bar | column | line | pie."""
    raw = str(b.get("chart_type") or b.get("chart") or b.get("kind") or "column").strip().lower()
    if "pie" in raw or "饼" in raw:
        return "pie"
    if "line" in raw or "折线" in raw or "趋势" in raw:
        return "line"
    if raw in ("bar", "barh", "horizontal", "条形"):
        return "bar"
    return "column"  # default: vertical clustered columns


def _normalize_series(b: dict) -> dict:
    """Coerce a chart's series into an ordered {name: [floats]} dict, shape-tolerant.

    Accepts {name: [vals]}, a list of {name, values} objects, or a bare list of
    numbers (one unnamed series).
    """
    raw = _first_key(b, _CHART_SERIES_KEYS)
    out: dict = {}
    if isinstance(raw, dict):
        for name, vals in raw.items():
            out[str(name)] = _floats(vals)
    elif isinstance(raw, list):
        if raw and all(isinstance(x, dict) for x in raw):
            for i, d in enumerate(raw):
                name = d.get("name") or d.get("label") or f"系列{i + 1}"
                vals = d.get("values") or d.get("data") or d.get("y") or []
                out[str(name)] = _floats(vals)
        elif raw and all(isinstance(x, (int, float)) for x in raw):
            out["系列1"] = _floats(raw)
    return out


def _normalize_block(b) -> dict | None:
    """Coerce one content block into a canonical {type, ...} dict, shape-tolerant.

    Recognized types: "paragraph" (成段正文), "bullets" (要点), "table" (表格).
    A bare string / number becomes a paragraph. An unknown-but-texty dict falls
    back to a paragraph using its first text-ish value. Returns None for blocks
    with no usable content so the caller can skip them.
    """
    if isinstance(b, (str, int, float)):
        text = str(b).strip()
        return {"type": "paragraph", "text": text} if text else None
    if not isinstance(b, dict):
        return None

    btype = str(b.get("type") or b.get("kind") or "").strip().lower()

    # Images: explicit type, or a src/url/path/prompt key. A `prompt` (no src)
    # means "generate this image" — resolved at render time if a backend exists.
    has_src = any(b.get(k) for k in _IMAGE_SRC_KEYS)
    prompt = b.get("prompt") or b.get("image_prompt")
    if btype in _IMAGE_TYPES or has_src or (prompt and not btype):
        src = ""
        for k in _IMAGE_SRC_KEYS:
            v = b.get(k)
            if isinstance(v, str) and v.strip():
                src = v.strip()
                break
        caption = ""
        for k in _IMAGE_CAPTION_KEYS:
            v = b.get(k)
            if isinstance(v, str) and v.strip():
                caption = v.strip()
                break
        if src or prompt:
            return {"type": "image", "src": src or None,
                    "prompt": str(prompt) if prompt else None, "caption": caption or None}
        if btype in _IMAGE_TYPES:
            return None  # explicit image with nothing usable → skip

    # Charts: explicit type, or duck-typed via a categories list (+ some series).
    cats = _first_key(b, _CHART_CAT_KEYS)
    if btype in _CHART_TYPES or (isinstance(cats, list) and bool(cats) and not btype):
        categories = [str(c) for c in cats] if isinstance(cats, list) else []
        series = _normalize_series(b)
        if categories and series:
            return {"type": "chart", "chart_type": _chart_kind(b),
                    "categories": categories, "series": series}
        if btype in _CHART_TYPES:
            return None  # explicit chart but no usable data → skip

    # Tables: explicit type, or duck-typed via row/header keys.
    rows = _first_key(b, _TABLE_ROWS_KEYS)
    headers = _first_key(b, _TABLE_HEADER_KEYS)
    if btype in ("table", "grid") or (rows is not None and btype not in ("paragraph", "bullets")):
        norm_headers = [str(h) for h in headers] if isinstance(headers, list) else []
        norm_rows = []
        if isinstance(rows, list):
            for r in rows:
                if isinstance(r, list):
                    norm_rows.append([str(c) for c in r])
                elif isinstance(r, dict):
                    # row as object → take values in header order if possible, else dict order.
                    if norm_headers and all(h in r for h in norm_headers):
                        norm_rows.append([str(r.get(h, "")) for h in norm_headers])
                    else:
                        norm_rows.append([str(v) for v in r.values()])
                else:
                    norm_rows.append([str(r)])
        if norm_headers or norm_rows:
            return {"type": "table", "headers": norm_headers, "rows": norm_rows}
        return None

    # Bullets: explicit type, or has a list of items under a bullet/child key.
    items = _first_key(b, _BULLET_CHILD_KEYS) or b.get("items")
    if btype in ("bullets", "bullet", "list", "ul") or (isinstance(items, list) and not btype):
        if isinstance(items, list) and items:
            return {"type": "bullets", "items": items}
        # explicitly a bullets block but empty → skip
        if btype:
            return None

    # Paragraph (default): pull the first text-ish value.
    text = ""
    for k in _PARA_TEXT_KEYS:
        v = b.get(k)
        if isinstance(v, str) and v.strip():
            text = v
            break
        if isinstance(v, (int, float)):
            text = str(v)
            break
    if not text:
        text = _bullet_text(b)  # last resort: reuse the broad bullet-text lookup
    return {"type": "paragraph", "text": text} if text else None


def _normalize_slide(s, idx: int) -> dict:
    """Coerce one slide entry into {title, blocks, notes}, tolerant of shape.

    Accepts a string (used as the title) or a dict using any of the common
    alternative key names. A missing title is non-fatal — it defaults to
    "第 N 页" so the deck still builds.

    `blocks` is an ordered list of canonical block dicts. A slide that only
    carries legacy `bullets` (no `blocks`) is normalized to a single bullets
    block tagged `_legacy` so build() can keep rendering it via the original
    content placeholder (zero visual regression).
    """
    if isinstance(s, str):
        return {"title": s, "blocks": [], "notes": None}
    if not isinstance(s, dict):
        return {"title": f"第 {idx} 页",
                "blocks": [{"type": "paragraph", "text": str(s)}], "notes": None}

    title = _first_key(s, _TITLE_KEYS) or f"第 {idx} 页"
    notes = _first_key(s, _NOTES_KEYS)

    raw_blocks = _first_key(s, _BLOCKS_KEYS)
    if isinstance(raw_blocks, list):
        blocks = [nb for nb in (_normalize_block(b) for b in raw_blocks) if nb]
        return {"title": str(title), "blocks": blocks, "notes": notes}

    # Legacy path: only bullets → single bullets block on the placeholder route.
    bullets = _first_key(s, _BULLETS_KEYS) or []
    if isinstance(bullets, str):
        # A single string body → split into lines as bullets.
        bullets = [ln for ln in bullets.splitlines() if ln.strip()] or [bullets]
    elif not isinstance(bullets, list):
        bullets = [str(bullets)]
    legacy = [{"type": "bullets", "items": bullets, "_legacy": True}] if bullets else []
    return {"title": str(title), "blocks": legacy, "notes": notes}


# ── block rendering (paragraph / bullets / table) ────────────────────────────

# Content region geometry (inches) for blocks-based slides, below the title.
_CONTENT_LEFT = 0.6
_CONTENT_TOP = 1.5
_CONTENT_WIDTH = 8.8
_CONTENT_BOTTOM = 7.0  # stop placing further blocks past this to avoid overflow
_BLOCK_GAP = 0.2       # vertical gap between blocks
_CHART_HEIGHT = 3.2    # default height (inches) for a chart block
_IMG_MAX_WIDTH = 5.0   # max width (inches) for an image block
_IMG_MAX_HEIGHT = 3.6  # max height (inches) for an image block
_EMU_PER_INCH = 914400
_BAND_HEIGHT = 1.35    # height (inches) of the colored title band

# Image extensions python-pptx can embed directly; anything else (e.g. .webp)
# is converted to PNG via Pillow before insertion.
_PPTX_NATIVE_EXT = (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".emf", ".wmf")

# Image-fetch limits. One slow/failed download must not blow the whole build's
# tool-timeout budget (which would kill the subprocess before prs.save(), losing
# the entire deck). Each fetch is capped, and a per-build cumulative budget stops
# further network fetches once exceeded — the deck still finishes with whatever
# images already succeeded; local files are unaffected.
_IMG_DOWNLOAD_TIMEOUT = 8   # seconds per single image download
_IMG_TIME_BUDGET = 20.0     # seconds total spent fetching/generating images per build
_img_time_spent = 0.0


def _reset_image_budget() -> None:
    global _img_time_spent
    _img_time_spent = 0.0


# ── themes (code-based color/font palettes, no external assets) ──────────────

# Palettes live in themes.py so the API server can serve the very same colors to
# the frontend's preview cards (one source of truth). Ensure this script's own
# directory is importable in BOTH contexts that load build.py: run as a
# subprocess (cwd = skill dir) and loaded by tests via importlib.
_HERE = os.path.dirname(os.path.abspath(__file__))
if _HERE not in sys.path:
    sys.path.insert(0, _HERE)
from themes import THEMES as _THEMES  # noqa: E402 — sys.path tweak must precede this import
_THEME_KEYS = ("title", "text", "background", "bg2", "font", "band", "band_text",
               "footer", "table_header", "table_header_text", "table_body")


def _hex(value):
    """Parse a hex color string (with/without '#', 3- or 6-digit) → RGBColor.

    Returns None on empty / invalid input so a bad palette value is skipped
    rather than crashing the whole build.
    """
    if not value:
        return None
    from pptx.dml.color import RGBColor

    s = str(value).lstrip("#").strip()
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    try:
        return RGBColor.from_string(s)
    except (ValueError, TypeError):
        return None


def _resolve_theme(spec: dict) -> dict:
    """Resolve spec["theme"] into a palette dict, shape-tolerant.

    Accepts a theme name (str) from _THEMES, or an inline object that overrides
    a named `base` (default "default"). Unknown names fall back to "default".
    """
    raw = spec.get("theme") or spec.get("style") or "default"
    if isinstance(raw, str):
        return dict(_THEMES.get(raw.strip().lower(), _THEMES["default"]))
    if isinstance(raw, dict):
        base = str(raw.get("base") or "default").strip().lower()
        merged = dict(_THEMES.get(base, _THEMES["default"]))
        for k in _THEME_KEYS:
            if raw.get(k):
                merged[k] = raw[k]
        return merged
    return dict(_THEMES["default"])


def _apply_text_theme(text_frame, theme: dict, color_key: str = "text") -> None:
    """Apply the theme's font name and a color to every paragraph in a frame."""
    color = _hex(theme.get(color_key))
    font = theme.get("font")
    for para in text_frame.paragraphs:
        if color is not None:
            para.font.color.rgb = color
        if font:
            para.font.name = font


def _apply_background(slide, theme: dict) -> None:
    """Fill the slide background — gradient if `bg2` set, else solid. Best-effort.

    No `background` → returns immediately (so the default theme stays blank).
    """
    color = _hex(theme.get("background"))
    if color is None:
        return
    color2 = _hex(theme.get("bg2"))
    try:
        fill = slide.background.fill
        if color2 is not None:
            try:
                fill.gradient()
                stops = fill.gradient_stops
                stops[0].position = 0.0
                stops[0].color.rgb = color
                stops[1].position = 1.0
                stops[1].color.rgb = color2
                try:
                    fill.gradient_angle = 90.0  # top → bottom
                except (NotImplementedError, ValueError):
                    pass
                return
            except Exception:  # noqa: BLE001 — gradient unsupported → fall back to solid
                fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color
    except Exception as e:  # noqa: BLE001 — background is cosmetic, never fatal
        print(f"[skip background: {type(e).__name__}: {e}]", file=sys.stderr)


def _send_to_back(slide, shape) -> None:
    """Move a shape to the bottom of the z-order (behind text/placeholders)."""
    try:
        sp_tree = slide.shapes._spTree
        el = shape._element
        sp_tree.remove(el)
        # Insert after <nvGrpSpPr> (0) and <grpSpPr> (1) so it is the bottom shape.
        sp_tree.insert(2, el)
    except Exception:  # noqa: BLE001 — z-order tweak is cosmetic, never fatal
        pass


def _no_border(shape) -> None:
    """Strip a decorative shape's outline and inherited shadow. Best-effort."""
    try:
        shape.line.fill.background()
    except Exception:  # noqa: BLE001
        pass
    try:
        shape.shadow.inherit = False
    except Exception:  # noqa: BLE001
        pass


def _add_header_band(slide, theme: dict, width_in: float) -> float:
    """Draw a full-width color band behind the title; returns its height (inches),
    or 0.0 if no band was drawn (so the caller can position the title in it)."""
    band = _hex(theme.get("band"))
    if band is None:
        return 0.0
    try:
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE

        shp = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(width_in), Inches(_BAND_HEIGHT))
        shp.fill.solid()
        shp.fill.fore_color.rgb = band
        _no_border(shp)
        _send_to_back(slide, shp)
        return _BAND_HEIGHT
    except Exception as e:  # noqa: BLE001 — decoration is cosmetic, never fatal
        print(f"[skip header band: {type(e).__name__}: {e}]", file=sys.stderr)
        return 0.0


def _add_footer(slide, theme: dict, page_no: int, width_in: float, height_in: float) -> None:
    """Draw a thin footer bar + page number in the bottom margin. Best-effort."""
    color = _hex(theme.get("footer") or theme.get("band"))
    if color is None:
        return
    try:
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN

        bar_h = 0.35
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(height_in - bar_h),
            Inches(width_in), Inches(bar_h))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        _no_border(bar)
        _send_to_back(slide, bar)

        box = slide.shapes.add_textbox(
            Inches(width_in - 1.2), Inches(height_in - bar_h), Inches(1.0), Inches(bar_h))
        tf = box.text_frame
        tf.text = str(page_no)
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.RIGHT
        para.font.size = Pt(10)
        num_color = _hex(theme.get("band_text")) or _hex("FFFFFF")
        if num_color is not None:
            para.font.color.rgb = num_color
    except Exception as e:  # noqa: BLE001 — footer is cosmetic, never fatal
        print(f"[skip footer: {type(e).__name__}: {e}]", file=sys.stderr)


def _decorate_cover(slide, theme: dict, width_in: float, height_in: float) -> None:
    """Decorate the title (cover) slide: a prominent top header band + a slim
    bottom accent. The top band is at least as tall as the bottom (a header
    narrower than the footer looks wrong). Centered title/subtitle placeholders
    are not overlapped. No-op when the theme has no band color.
    """
    band = _hex(theme.get("band"))
    if band is None:
        return
    try:
        from pptx.util import Inches
        from pptx.enum.shapes import MSO_SHAPE

        # Prominent header band across the top.
        top = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(width_in), Inches(1.0))
        top.fill.solid()
        top.fill.fore_color.rgb = band
        _no_border(top)
        _send_to_back(slide, top)

        # Slim accent bar along the bottom (thinner than the header).
        bottom = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(height_in - 0.45),
            Inches(width_in), Inches(0.45))
        bottom.fill.solid()
        bottom.fill.fore_color.rgb = band
        _no_border(bottom)
        _send_to_back(slide, bottom)
    except Exception as e:  # noqa: BLE001 — decoration is cosmetic, never fatal
        print(f"[skip cover decor: {type(e).__name__}: {e}]", file=sys.stderr)


# ── layout / title helpers (work with arbitrary templates, not fixed indices) ─

def _pick_layout(prs, names, fallback_idx: int):
    """Pick a slide layout by name (case-insensitive substring), in priority order.

    Falls back to a positional index (then the first layout) when no name
    matches — so an arbitrary template without the usual layout names still
    yields a usable layout.
    """
    for want in names:
        w = want.lower()
        for layout in prs.slide_layouts:
            if w in (layout.name or "").lower():
                return layout
    layouts = list(prs.slide_layouts)
    if 0 <= fallback_idx < len(layouts):
        return layouts[fallback_idx]
    return layouts[0]


def _title_font_size(text: str) -> int:
    """Pick a title font size (pt) that shrinks for longer titles so they stay
    within ~2 lines inside the header band (CJK chars count as wider)."""
    width = sum(2 if ord(c) > 0x2E80 else 1 for c in str(text))  # CJK ≈ 2 cells
    if width <= 28:
        return 32
    if width <= 44:
        return 28
    if width <= 64:
        return 22
    return 18


def _enable_shrink_to_fit(text_frame) -> None:
    """Tag the text frame so PowerPoint shrinks the text if it still overflows
    (belt-and-suspenders on top of the length-based font sizing). Best-effort."""
    try:
        from pptx.oxml.ns import qn

        body_pr = text_frame._txBody.bodyPr
        for tag in ("a:noAutofit", "a:normAutofit", "a:spAutoFit"):
            for el in body_pr.findall(qn(tag)):
                body_pr.remove(el)
        body_pr.append(body_pr.makeelement(qn("a:normAutofit"), {}))
    except Exception:  # noqa: BLE001 — autofit hint is cosmetic, never fatal
        pass


def _set_title(slide, text: str, theme: dict, color_key: str = "title",
               band_box: tuple | None = None) -> None:
    """Set the slide title, using the title placeholder if the layout has one,
    else a manually-added textbox (for templates whose layout lacks a title).

    `color_key` selects which theme color to paint the title with — "band_text"
    when the title sits on a colored header band, "title" otherwise.
    `band_box` = (left, top, width, height) in inches: when given, the title is
    repositioned to fill the band, vertically centered, word-wrapped, and sized
    down for long text so it does not overflow the band."""
    from pptx.util import Inches, Pt
    from pptx.enum.text import MSO_ANCHOR

    title_ph = slide.shapes.title
    if title_ph is not None:
        title_ph.text = str(text)
        target = title_ph
    else:
        target = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        target.text_frame.text = str(text)

    tf = target.text_frame
    para = tf.paragraphs[0]
    para.font.bold = True

    if band_box is not None:
        left, top, width, height = band_box
        target.left, target.top = Inches(left), Inches(top)
        target.width, target.height = Inches(width), Inches(height)
        tf.word_wrap = True
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE  # vertical centering within the band
            tf.margin_top = Inches(0.04)
            tf.margin_bottom = Inches(0.04)
        except Exception:  # noqa: BLE001
            pass
        para.font.size = Pt(_title_font_size(str(text)))
        _enable_shrink_to_fit(tf)
    elif title_ph is None:
        para.font.size = Pt(28)  # manual-textbox fallback without a band

    _apply_text_theme(tf, theme, color_key)


def _content_placeholder(slide):
    """Return the body/content placeholder (idx 1) of a slide, or None."""
    try:
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 1:
                return ph
    except Exception:  # noqa: BLE001
        return None
    return None


def _add_paragraph(text_frame, text: str) -> None:
    """Write a (possibly multi-line) prose paragraph into a text frame."""
    text_frame.word_wrap = True
    first = True
    for line in str(text).split("\n"):
        para = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
        para.text = line
        first = False


def _est_height(block: dict, width_in: float = _CONTENT_WIDTH) -> float:
    """Rough height (inches) a block will occupy, for the flow layout.

    `width_in` is the column width the block is laid out in; a narrower column
    wraps text into more lines, so the per-line char budget scales with it.
    """
    btype = block.get("type")
    if btype == "paragraph":
        text = str(block.get("text", ""))
        # ~45 chars per line at the full content width / 16pt; scale to the column.
        cpl = max(12, int(45 * width_in / _CONTENT_WIDTH))
        lines = sum(max(1, (len(ln) // cpl) + 1) for ln in text.split("\n"))
        return max(0.4, lines * 0.3)
    if btype == "bullets":
        # count top-level + nested items shallowly
        items = block.get("items") or []
        return max(0.4, len(items) * 0.35)
    if btype == "table":
        return max(0.4, (len(block.get("rows") or []) + 1) * 0.35)
    if btype == "chart":
        return _CHART_HEIGHT
    if btype == "image":
        # actual height is known only after the picture loads; reserve a default
        # (plus the caption line so the fit check leaves room for it).
        return _IMG_MAX_HEIGHT + (0.35 if block.get("caption") else 0.0)
    return 0.4


def _add_table(slide, left, top, width, headers, rows, theme: dict) -> None:
    """Render a table; header row bold on a themed fill. Caller positions it."""
    from pptx.util import Inches, Pt

    body = list(rows or [])
    ncols = max([len(headers or [])] + [len(r) for r in body] or [0]) or 1
    nrows = len(body) + (1 if headers else 0)
    if nrows == 0:
        return
    header_fill = _hex(theme.get("table_header")) or _hex("1F4E79")
    header_text = _hex(theme.get("table_header_text")) or _hex("FFFFFF")
    body_text = _hex(theme.get("text"))
    body_fill = _hex(theme.get("table_body"))
    font = theme.get("font")

    height = Inches(_est_height({"type": "table", "rows": body}))
    gfx = slide.shapes.add_table(nrows, ncols, Inches(left), Inches(top),
                                 Inches(width), height)
    table = gfx.table

    r = 0
    if headers:
        for c in range(ncols):
            cell = table.cell(0, c)
            cell.text = headers[c] if c < len(headers) else ""
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
            para = cell.text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(14)
            para.font.color.rgb = header_text
            if font:
                para.font.name = font
        r = 1
    for row in body:
        for c in range(ncols):
            cell = table.cell(r, c)
            cell.text = row[c] if c < len(row) else ""
            # Paint an explicit cell fill so the themed body text isn't drawn onto
            # python-pptx's default (light, banded) table style — on dark themes
            # that left near-white text on near-white cells, unreadable.
            if body_fill is not None:
                cell.fill.solid()
                cell.fill.fore_color.rgb = body_fill
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            if body_text is not None:
                para.font.color.rgb = body_text
            if font:
                para.font.name = font
        r += 1


_CHART_XL = {
    "column": "COLUMN_CLUSTERED",
    "bar": "BAR_CLUSTERED",
    "line": "LINE_MARKERS",
    "pie": "PIE",
}


def _add_chart(slide, left, top, width, block, theme: dict) -> float:
    """Render a native (editable) chart. Returns the height consumed (inches)."""
    from pptx.util import Inches
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    categories = block.get("categories") or []
    series = block.get("series") or {}
    if not categories or not series:
        return 0.0
    n = len(categories)
    data = CategoryChartData()
    data.categories = categories
    kind = block.get("chart_type", "column")
    if kind == "pie":
        # Pie charts plot a single series; use the first.
        name, vals = next(iter(series.items()))
        data.add_series(name, (list(vals) + [0.0] * n)[:n])
    else:
        for name, vals in series.items():
            data.add_series(name, (list(vals) + [0.0] * n)[:n])

    xl = getattr(XL_CHART_TYPE, _CHART_XL.get(kind, "COLUMN_CLUSTERED"))
    height = Inches(_CHART_HEIGHT)
    gf = slide.shapes.add_chart(xl, Inches(left), Inches(top), Inches(width), height, data)
    chart = gf.chart
    chart.has_legend = kind == "pie" or len(series) > 1
    if chart.has_legend:
        chart.legend.include_in_layout = False
    return _CHART_HEIGHT


def _download_image(url: str) -> str | None:
    """Download an http(s) image into generated/ and return the local path, or None."""
    if not url.lower().startswith(("http://", "https://")):
        return None
    try:
        req = urllib.request.Request(url, headers={"User-Agent": "ppt-build/1.0"})
        with urllib.request.urlopen(req, timeout=_IMG_DOWNLOAD_TIMEOUT) as resp:  # noqa: S310 — scheme checked above
            ctype = resp.headers.get("Content-Type", "")
            data = resp.read(8 * 1024 * 1024)  # cap at 8 MB
    except Exception as e:  # noqa: BLE001 — network is best-effort; skip on any failure
        print(f"[image download failed: {type(e).__name__}: {e}]", file=sys.stderr)
        return None
    if not data:
        return None
    ext = ".png"
    for cand, suffix in ((".jpeg", ".jpg"), (".jpg", ".jpg"), (".png", ".png"),
                         (".gif", ".gif"), (".webp", ".webp")):
        if cand in ctype.lower() or url.lower().split("?")[0].endswith(cand):
            ext = suffix
            break
    out_dir = _generated_dir()
    os.makedirs(out_dir, exist_ok=True)
    path = os.path.join(out_dir, f"img-{uuid.uuid4().hex[:8]}{ext}")
    with open(path, "wb") as f:
        f.write(data)
    return path


def _generate_image(prompt: str) -> str | None:
    """Best-effort text-to-image via DashScope (Tongyi Wanxiang). Returns a local
    path or None — missing key / any error degrades silently (no key needed to test
    the rest of the pipeline)."""
    if not os.getenv("DASHSCOPE_API_KEY"):
        return None
    try:
        from dashscope import ImageSynthesis

        model = os.getenv("IMAGE_MODEL", "wanx2.1-t2i-turbo")
        rsp = ImageSynthesis.call(model=model, prompt=prompt, n=1, size="1024*1024")
        if getattr(rsp, "status_code", None) == 200 and rsp.output and rsp.output.results:
            return _download_image(rsp.output.results[0].url)
    except Exception as e:  # noqa: BLE001 — generation is optional
        print(f"[image generation failed: {type(e).__name__}: {e}]", file=sys.stderr)
    return None


def _resolve_image(src: str | None, prompt: str | None) -> str | None:
    """Turn an image block's src/prompt into a usable local file path, or None.

    src may be an http(s) URL (downloaded), an absolute/cwd-relative local file,
    or a bare filename present in generated/. Falls back to AI generation when
    only a prompt is given.

    Network fetches and generation draw on a per-build time budget
    (`_IMG_TIME_BUDGET`); once exhausted they are skipped so a few slow/failed
    images never stall the build into a tool-timeout kill. Local files always
    work regardless of the budget.
    """
    global _img_time_spent
    if src:
        if src.lower().startswith(("http://", "https://")):
            if _img_time_spent >= _IMG_TIME_BUDGET:
                print(f"[image budget spent, skipping download {src}]", file=sys.stderr)
            else:
                t0 = time.monotonic()
                local = _download_image(src)
                _img_time_spent += time.monotonic() - t0
                if local:
                    return local
        elif os.path.isfile(src):
            return src
        else:
            gen = os.path.join(_generated_dir(), os.path.basename(src))
            if os.path.isfile(gen):
                return gen
    if prompt:
        if _img_time_spent >= _IMG_TIME_BUDGET:
            print("[image budget spent, skipping generation]", file=sys.stderr)
            return None
        t0 = time.monotonic()
        out = _generate_image(prompt)
        _img_time_spent += time.monotonic() - t0
        return out
    return None


def _ensure_supported_image(path: str) -> str:
    """Convert an image python-pptx cannot embed (e.g. .webp) to PNG via Pillow.

    Returns a path python-pptx accepts: the original for native formats, a
    converted .png otherwise, or the original unchanged if conversion fails
    (in which case add_picture raises and the block is skipped cleanly)."""
    ext = os.path.splitext(path)[1].lower()
    if ext in _PPTX_NATIVE_EXT:
        return path
    try:
        from PIL import Image

        out = os.path.splitext(path)[0] + ".png"
        with Image.open(path) as im:
            im.convert("RGBA").save(out, "PNG")
        return out
    except Exception as e:  # noqa: BLE001 — fall back to the original path
        print(f"[image convert failed {ext!r}: {type(e).__name__}: {e}]", file=sys.stderr)
        return path


def _add_image(slide, left, top, avail_w, block, theme: dict,
               max_height: float = _IMG_MAX_HEIGHT, region_h: float | None = None) -> float:
    """Render an image (best-fit into the available box) + optional caption.

    The picture is scaled to fit *within* `avail_w` × `max_height` preserving its
    aspect ratio (so it is as large as the space allows — minimal surrounding
    whitespace), then centered horizontally in `avail_w`. When `region_h` is
    given, the image+caption block is also centered vertically inside that region
    (used by the two-column layout so the picture sits beside the text instead of
    hugging the top with empty space below). Returns the height consumed (inches).
    """
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    path = _resolve_image(block.get("src"), block.get("prompt"))
    if not path:
        return 0.0
    cap_reserve = 0.35 if block.get("caption") else 0.0
    avail_h = max(0.6, max_height - cap_reserve)
    path = _ensure_supported_image(path)
    # Add at native size, then scale to best-fit the box preserving aspect ratio.
    pic = slide.shapes.add_picture(path, Inches(left), Inches(top))
    nat_w = pic.width / _EMU_PER_INCH
    nat_h = pic.height / _EMU_PER_INCH
    if nat_w <= 0 or nat_h <= 0:
        return 0.0
    scale = min(avail_w / nat_w, avail_h / nat_h)
    w, h = nat_w * scale, nat_h * scale
    total = h + cap_reserve
    # Vertically center the whole block within region_h (if any), else sit at top.
    y = top + (region_h - total) / 2.0 if (region_h and region_h > total) else top
    x = left + max(0.0, (avail_w - w) / 2.0)  # horizontally center within avail_w
    pic.width, pic.height = Inches(w), Inches(h)
    pic.left, pic.top = Inches(x), Inches(y)
    if block.get("caption"):
        cap = slide.shapes.add_textbox(Inches(x), Inches(y + h + 0.05), Inches(w), Inches(0.3))
        cap.text_frame.word_wrap = True
        cap.text_frame.text = str(block["caption"])
        para = cap.text_frame.paragraphs[0]
        para.font.size = Pt(11)
        para.font.italic = True
        para.alignment = PP_ALIGN.CENTER
        _apply_text_theme(cap.text_frame, theme, "text")
    return total


def _render_two_column(slide, text_blocks, image_blocks, theme: dict,
                       bottom: float = _CONTENT_BOTTOM) -> bool:
    """Side-by-side layout: text in a left column, image(s) stacked on the right.

    Returns True only if at least one image was actually placed — the caller
    falls back to the vertical flow otherwise (e.g. an image that failed to
    download). This is what lets a text-heavy slide still show its picture: the
    vertical flow would drop an image appended after a full page of text.
    """
    col_gap = 0.4
    right_w = 4.0
    left_w = _CONTENT_WIDTH - right_w - col_gap
    left_x = _CONTENT_LEFT
    right_x = _CONTENT_LEFT + left_w + col_gap
    col_h = bottom - _CONTENT_TOP

    # Right column first: place the image(s), centered in the column both ways so
    # they sit beside the text rather than hugging the top with whitespace below.
    placed = False
    if len(image_blocks) == 1:
        try:
            placed = _add_image(slide, right_x, _CONTENT_TOP, right_w, image_blocks[0],
                                theme, max_height=col_h, region_h=col_h) > 0
        except Exception as e:  # noqa: BLE001 — one bad image must not kill the deck
            print(f"[skip image: {type(e).__name__}: {e}]", file=sys.stderr)
    else:
        img_top = _CONTENT_TOP
        for img in image_blocks:
            if img_top >= bottom - 0.6:
                break
            try:
                real = _add_image(slide, right_x, img_top, right_w, img, theme,
                                  max_height=bottom - img_top)
            except Exception as e:  # noqa: BLE001
                print(f"[skip image: {type(e).__name__}: {e}]", file=sys.stderr)
                real = 0.0
            if real > 0:
                placed = True
                img_top += real + _BLOCK_GAP
    if not placed:
        return False

    # Left column: ALL text in one word-wrapped, shrink-to-fit textbox spanning the
    # full column height. A single box (vs one per block) makes vertical overlap
    # impossible, word-wrap keeps long lines inside the column (never under the
    # image), and shrink-to-fit scales the text down if it is still too tall.
    _render_text_box(slide, text_blocks, theme, left_x, _CONTENT_TOP, left_w, col_h)
    return True


def _render_text_box(slide, text_blocks, theme: dict, left, top, width, height,
                     base_size: int = 15) -> None:
    """Render paragraph/bullets blocks into a single word-wrapped textbox."""
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for bi, block in enumerate(text_blocks):
        if not first:
            # blank spacer line between blocks
            tf.add_paragraph()
        if block.get("type") == "bullets":
            for item in (block.get("items") or []):
                txt = _bullet_text(item) if isinstance(item, dict) else str(item)
                if not txt:
                    continue
                para = tf.paragraphs[0] if first else tf.add_paragraph()
                para.text = "• " + txt  # manual bullet glyph (a plain textbox has no list style)
                para.font.size = Pt(base_size)
                first = False
        else:
            for line in str(block.get("text", "")).split("\n"):
                para = tf.paragraphs[0] if first else tf.add_paragraph()
                para.text = line
                para.font.size = Pt(base_size)
                first = False
    _apply_text_theme(tf, theme, "text")
    _enable_shrink_to_fit(tf)


def _render_blocks(slide, blocks, theme: dict, bottom: float = _CONTENT_BOTTOM) -> None:
    """Lay a slide's blocks out in the content region.

    When a slide mixes an image with text (and has no full-width table/chart),
    use a two-column text|image layout so the image is never crowded out. Every
    other shape mix flows top-to-bottom via `_render_blocks_flow`.
    """
    images = [b for b in blocks if b.get("type") == "image"]
    texts = [b for b in blocks if b.get("type") in ("paragraph", "bullets")]
    wide = [b for b in blocks if b.get("type") in ("table", "chart")]
    if images and texts and not wide:
        if _render_two_column(slide, texts, images, theme, bottom):
            return
    _render_blocks_flow(slide, blocks, theme, bottom)


def _render_blocks_flow(slide, blocks, theme: dict, bottom: float = _CONTENT_BOTTOM) -> None:
    """Lay blocks out top-to-bottom as independent shapes in the content region.

    `bottom` is the lowest inch the content may reach (kept above the footer by
    the caller). A block that would extend past it is dropped rather than drawn
    so content never overflows into the footer / off the slide.
    """
    from pptx.util import Inches, Pt

    top = _CONTENT_TOP
    for block in blocks:
        if top >= bottom:
            break  # ran out of vertical room — drop remaining blocks
        btype = block.get("type")
        try:
            est = _est_height(block)
            # Skip a block that won't fit above the footer (prevents overflow).
            # Always allow the very first block so a page is never left empty.
            # Images are exempt: rather than drop them, they shrink to the space
            # that remains (below) — a small picture beats no picture.
            if btype != "image" and top > _CONTENT_TOP and top + est > bottom:
                break
            if btype == "table":
                _add_table(slide, _CONTENT_LEFT, top, _CONTENT_WIDTH,
                           block.get("headers"), block.get("rows"), theme)
            elif btype == "chart":
                est = _add_chart(slide, _CONTENT_LEFT, top, _CONTENT_WIDTH, block, theme) or 0.0
            elif btype == "image":
                # Fit the image into the full content width and the room left,
                # centered — as large as the space allows (less whitespace), and
                # never overflowing the footer.
                real = _add_image(slide, _CONTENT_LEFT, top, _CONTENT_WIDTH, block, theme,
                                  max_height=max(0.8, bottom - top))
                est = real or 0.0  # unresolved image → consume nothing, skip cleanly
            elif btype == "bullets":
                box = slide.shapes.add_textbox(
                    Inches(_CONTENT_LEFT), Inches(top), Inches(_CONTENT_WIDTH), Inches(est))
                box.text_frame.word_wrap = True  # keep long bullet lines inside the box
                _add_bullets(box.text_frame, block.get("items") or [])
                for para in box.text_frame.paragraphs:
                    para.font.size = Pt(16)
                _apply_text_theme(box.text_frame, theme, "text")
            else:  # paragraph
                box = slide.shapes.add_textbox(
                    Inches(_CONTENT_LEFT), Inches(top), Inches(_CONTENT_WIDTH), Inches(est))
                _add_paragraph(box.text_frame, block.get("text", ""))
                for para in box.text_frame.paragraphs:
                    para.font.size = Pt(16)
                _apply_text_theme(box.text_frame, theme, "text")
            if est > 0:
                top += est + _BLOCK_GAP
        except Exception as e:  # noqa: BLE001 — one bad block must not kill the deck
            print(f"[skip block type={btype}: {type(e).__name__}: {e}]", file=sys.stderr)


def _templates_dir() -> str:
    """Directory of bundled .pptx template files, alongside this script."""
    here = os.path.dirname(os.path.abspath(__file__))
    return os.path.join(here, "templates")


def _open_presentation(spec: dict):
    """Open the base Presentation: a bundled template if spec asks for a valid
    one, else python-pptx's built-in default. Always falls back gracefully."""
    from pptx import Presentation

    tmpl = spec.get("template")
    if isinstance(tmpl, str) and tmpl.strip():
        base = os.path.abspath(_templates_dir())
        name = os.path.basename(tmpl.strip())
        if not name.lower().endswith(".pptx"):
            name += ".pptx"
        path = os.path.abspath(os.path.join(base, name))
        # Path guard: must resolve inside templates/ and exist.
        if path.startswith(base + os.sep) and os.path.isfile(path):
            try:
                return Presentation(path)
            except Exception as e:  # noqa: BLE001 — fall back to default on any error
                print(f"[template load failed {name!r}: {type(e).__name__}: {e}]",
                      file=sys.stderr)
        else:
            print(f"[template not found {name!r} — using default]", file=sys.stderr)
    return Presentation()


def build(spec: dict) -> str:
    from pptx.util import Pt

    raw_slides = spec.get("slides") or []
    if not isinstance(raw_slides, list):
        raise ValueError("'slides' must be a list")
    slides = [_normalize_slide(s, i + 1) for i, s in enumerate(raw_slides)]
    # Title is optional: fall back to the first slide's heading, else a default,
    # so a slides-only spec still builds a usable deck.
    title = _first_key(spec, _TITLE_KEYS)
    if not title:
        title = slides[0]["title"] if slides else "演示文稿"

    prs = _open_presentation(spec)
    theme = _resolve_theme(spec)
    _reset_image_budget()  # per-build image-fetch time budget (see _resolve_image)
    width_in = prs.slide_width / _EMU_PER_INCH
    height_in = prs.slide_height / _EMU_PER_INCH
    # Keep block content above the footer bar (0.35") with a small gap.
    content_bottom = height_in - 0.55

    # ── Title slide (matched by layout name, not a fixed index) ────────────────
    title_layout = _pick_layout(prs, ["title slide", "title and subtitle"], 0)
    title_slide = prs.slides.add_slide(title_layout)
    _apply_background(title_slide, theme)
    _decorate_cover(title_slide, theme, width_in, height_in)
    _set_title(title_slide, str(title), theme)  # cover title uses the normal color
    subtitle = spec.get("subtitle")
    if subtitle:
        try:
            if len(title_slide.placeholders) > 1:
                ph = title_slide.placeholders[1]
                ph.text = str(subtitle)
                _apply_text_theme(ph.text_frame, theme, "text")
        except Exception as e:  # noqa: BLE001 — subtitle is optional, never fatal
            print(f"[skip subtitle: {type(e).__name__}: {e}]", file=sys.stderr)

    # ── Content slides ────────────────────────────────────────────────────────
    for page_no, s in enumerate(slides, start=1):
        blocks = s["blocks"]
        legacy = bool(blocks) and all(b.get("_legacy") for b in blocks)

        if legacy:
            # Single legacy bullets block → keep the Title+Content placeholder
            # path for zero visual regression (matched by layout name).
            layout = _pick_layout(prs, ["title and content", "content"], 1)
            slide = prs.slides.add_slide(layout)
            _apply_background(slide, theme)
            band_h = _add_header_band(slide, theme, width_in)  # behind the title
            band_box = (0.5, 0.0, max(1.0, width_in - 1.0), band_h) if band_h else None
            _set_title(slide, s["title"], theme,
                       "band_text" if band_h else "title", band_box)
            body_ph = _content_placeholder(slide)
            if body_ph is not None:
                body = body_ph.text_frame
                _add_bullets(body, blocks[0]["items"])
                for para in body.paragraphs:
                    para.font.size = Pt(18)
                _apply_text_theme(body, theme, "text")
            else:
                # Template's layout lacks a body placeholder → flow as blocks.
                _render_blocks(slide, blocks, theme, content_bottom)
        else:
            # Title Only layout keeps the themed title; blocks flow below it.
            layout = _pick_layout(prs, ["title only", "blank"], 5)
            slide = prs.slides.add_slide(layout)
            _apply_background(slide, theme)
            band_h = _add_header_band(slide, theme, width_in)
            band_box = (0.5, 0.0, max(1.0, width_in - 1.0), band_h) if band_h else None
            _set_title(slide, s["title"], theme,
                       "band_text" if band_h else "title", band_box)
            if blocks:
                _render_blocks(slide, blocks, theme, content_bottom)

        _add_footer(slide, theme, page_no, width_in, height_in)

        if s["notes"]:
            slide.notes_slide.notes_text_frame.text = str(s["notes"])

    # Save into the repo's shared generated/ dir so the web server can serve it
    # for download. Append a suffix derived from the spec's CONTENT (not a random
    # token) so concurrent users with different decks never overwrite each other,
    # while the same spec always maps to the same filename. That idempotence is
    # what stops the model from looping: re-running an identical build returns the
    # exact same name it saw last time, instead of a fresh random name that reads
    # as "the save didn't take" and triggers yet another retry. Keep only the
    # basename's stem (no path separators) so the result is always one servable file.
    output = spec.get("output") or "deck.pptx"
    stem = os.path.splitext(os.path.basename(str(output)))[0].strip() or "deck"
    digest = hashlib.sha1(
        json.dumps(spec, ensure_ascii=False, sort_keys=True).encode("utf-8")
    ).hexdigest()[:8]
    name = f"{stem}-{digest}.pptx"
    out_dir = _generated_dir()
    os.makedirs(out_dir, exist_ok=True)
    path = os.path.join(out_dir, name)
    prs.save(path)
    return path


def _parse_spec(raw: str) -> dict:
    """Parse the CLI arg into a spec dict, tolerating double-encoded JSON.

    LLM/tool serialization frequently JSON-encodes the spec one extra time, so
    a single json.loads yields a *string* rather than the object. Re-parse while
    the result is still a string before giving up.

    `strict=False` tolerates literal control characters (real newlines/tabs)
    inside string values: the model writes `\\n` in its spec, but the tool-call
    argument is JSON-decoded once before reaching this script, turning those into
    real newlines — which strict JSON parsing would reject as
    "Invalid control character", failing the whole build.
    """
    value = json.loads(raw, strict=False)
    for _ in range(3):
        if isinstance(value, str):
            value = json.loads(value, strict=False)
        else:
            break
    # Unwrap a spec object that the model wrapped in a 1-element array, e.g.
    # `[{"title": ..., "slides": [...]}]`. Without this the whole spec would be
    # mistaken for a single slide and all real slides would be lost.
    if (
        isinstance(value, list)
        and len(value) == 1
        and isinstance(value[0], dict)
        and ("slides" in value[0] or "subtitle" in value[0])
    ):
        value = value[0]
    # Otherwise tolerate a bare list of slides (a common shape the model emits):
    # wrap it into a proper spec so the deck still builds.
    if isinstance(value, list):
        value = {"slides": value}
    if not isinstance(value, dict):
        raise ValueError(
            f"spec must be a JSON object or a list of slides, got {type(value).__name__}"
        )
    return value


def main() -> None:
    # build.py takes a SINGLE JSON spec, but the model sometimes splits it across
    # multiple script_args elements — each of which run_skill_script forwards as a
    # separate argv. Reading only argv[1] then parses a truncated prefix and fails
    # with e.g. "Expecting ',' delimiter" exactly at the split boundary. Rejoin all
    # args (they were passed with no separator) so a split spec is reconstructed.
    raw = "".join(sys.argv[1:])
    if not raw.strip():
        print(json.dumps({"ok": False, "error": "usage: build.py '<json-spec>'"}))
        sys.exit(1)
    try:
        spec = _parse_spec(raw)
    except (json.JSONDecodeError, ValueError) as e:
        print(json.dumps({"ok": False, "error": f"invalid JSON spec: {e}"}, ensure_ascii=False))
        sys.exit(1)
    try:
        path = build(spec)
    except Exception as e:  # noqa: BLE001 — report any failure as structured JSON
        print(json.dumps({"ok": False, "error": f"{type(e).__name__}: {e}"}, ensure_ascii=False))
        sys.exit(1)
    n = len(spec.get("slides") or [])
    print(json.dumps(
        {
            "ok": True,
            "path": path,
            "filename": os.path.basename(path),
            "slides": n,
            # The stop signal lives in the result itself (not just SKILL.md) so
            # the model reliably sees it and does not loop calling build.py.
            "message": (
                f"演示文稿已生成完毕，共 {n} 页，保存为 {os.path.basename(path)}。"
                "build.py 每次都会覆盖整份文件（不是追加），此文件已包含完整内容。"
                "不要再次调用 build.py——请直接把结果告知用户。"
            ),
        },
        ensure_ascii=False,
    ))


if __name__ == "__main__":
    main()
